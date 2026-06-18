"""nature-paper2ppt compiler adapter — paper outline -> a REAL .pptx slide deck.
 Entrypoint: render(...) per
../ADAPTER-CONTRACT.md (compiler kind), backing `render slides`.

Pipeline (degrades gracefully, always records an artifact):
  1. PRIMARY: python-pptx -> a real 16:9 `.pptx` deck. One slide per outline section (title, bullets,
     speaker notes, embedded hero figure + source label). Honors the `--ref` output path.
       - Figure assets are gathered by:
           a. PyMuPDF/`fitz` extraction from a source paper PDF (params/JSON `paper_pdf`), then
           b. figures already on disk under `runs/<quest_id>/figures/` (or JSON `figures_dir`).
       - Only "story" figures are selected (hero panel per slide), per the Lean Operating Mode.
       - Writes `asset_manifest.md` (figure traceability) alongside the deck.
       - QA: reopens the produced .pptx, counts slides + embedded media + notes, writes `qa_report.md`.
  2. FALLBACK: if python-pptx is unavailable at runtime, emit a self-contained HTML deck at the same ref
     (mirrors how paper-latex degrades). Figure extraction degrades to "skip" when `fitz` is absent.

Input JSON: {"title": "...", "author": "...",
             "paper_pdf": "<optional path>", "figures_dir": "<optional path>",
             "slides": [{"title": "...", "bullets": [...], "notes": "...",
                         "figure": "<optional asset path or basename>",
                         "caption": "...", "takeaway": "...", "source": "Fig. 2b, Nature, 2024"}]}
Input MD:   "# Deck title" then "## Slide title" sections with "- bullet" lines.

params (all optional): {"title": str, "paper_pdf": "<path>", "figures_dir": "<path>", "lang": "cn|en"}
out_path / --ref: target `.pptx` path (or `.html` in fallback). Side files land in its parent dir.
"""
from __future__ import annotations
import html
import json
import re
import zipfile
from pathlib import Path


# ───────────────────────────── input parsing ─────────────────────────────
def _parse(input_path):
    p = Path(input_path)
    if not p.exists():
        raise FileNotFoundError(f"nature-paper2ppt: input not found: {input_path}")
    if p.suffix.lower() == ".json":
        d = json.loads(p.read_text(encoding="utf-8"))
        return (d.get("title", "Presentation"), d.get("author", ""), d.get("slides", []) or [],
                d.get("paper_pdf"), d.get("figures_dir"))
    # Markdown: "# Deck title", "## Slide title", "- bullet", "> note" (speaker note),
    # "![](fig.png)" (figure asset). Conservative parse to mirror the JSON schema.
    title, slides, cur = "Presentation", [], None
    for line in p.read_text(encoding="utf-8").splitlines():
        s = line.strip()
        if s.startswith("## "):
            cur = {"title": s[3:].strip(), "bullets": []}
            slides.append(cur)
        elif s.startswith("# "):
            title = s[2:].strip()
        elif s.startswith(("- ", "* ")) and cur is not None:
            cur["bullets"].append(s[2:].strip())
        elif s.startswith("> ") and cur is not None:
            cur["notes"] = (cur.get("notes", "") + " " + s[2:].strip()).strip()
        elif cur is not None:
            m = re.match(r"!\[[^\]]*\]\(([^)]+)\)", s)
            if m:
                cur["figure"] = m.group(1).strip()
    return title, "", slides, None, None


# ───────────────────────────── figure assets ─────────────────────────────
_IMG_EXT = (".png", ".jpg", ".jpeg", ".gif", ".bmp", ".tif", ".tiff", ".webp")


def _figures_root(quest_id, figures_dir):
    """Resolve the on-disk figure pool. Explicit figures_dir wins; else runs/<quest>/figures/ then outputs/."""
    cands = []
    if figures_dir:
        cands.append(Path(figures_dir))
    # adapter sits at execplan/packs/nature-paper2ppt/adapter.py -> loop root is 3 parents up.
    loop_root = Path(__file__).resolve().parents[3]
    if quest_id:
        cands.append(loop_root / "runs" / quest_id / "figures")
    cands.append(loop_root / "runs" / "figures")
    cands.append(loop_root / "outputs" / "figures")
    for c in cands:
        if c and c.is_dir():
            return c
    return None


def _extract_pdf_figures(paper_pdf, asset_dir):
    """Extract embedded images from a paper PDF via PyMuPDF (fitz) into asset_dir.

    Returns (list[ {path, page, xref, w, h} ], method_str). Degrades to ([], 'fitz-absent') when PyMuPDF
    is not installed and ([], 'no-pdf') when there is no readable PDF. Only reasonably large images are
    kept (skips icons/logos/rules) to honor 'select only story figures'.
    """
    if not paper_pdf:
        return [], "no-pdf"
    pdf = Path(paper_pdf)
    if not pdf.exists():
        return [], "no-pdf"
    try:
        import fitz  # PyMuPDF
    except Exception:
        return [], "fitz-absent"
    try:
        from PIL import Image  # noqa: F401  (used only to validate; PyMuPDF writes the file)
    except Exception:
        pass
    asset_dir.mkdir(parents=True, exist_ok=True)
    out, seen = [], set()
    try:
        doc = fitz.open(str(pdf))
    except Exception:
        return [], "pdf-open-failed"
    try:
        for pno in range(doc.page_count):
            page = doc[pno]
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen:
                    continue
                seen.add(xref)
                try:
                    pix = fitz.Pixmap(doc, xref)
                except Exception:
                    continue
                # Skip tiny assets (logos, rules, bullets): require a meaningful footprint.
                if pix.width < 120 or pix.height < 90 or (pix.width * pix.height) < 30000:
                    continue
                try:
                    if pix.n - pix.alpha >= 4:  # CMYK -> RGB
                        pix = fitz.Pixmap(fitz.csRGB, pix)
                    fname = f"pdffig_p{pno + 1}_x{xref}.png"
                    fpath = asset_dir / fname
                    pix.save(str(fpath))
                    out.append({"path": fpath, "page": pno + 1, "xref": xref,
                                "w": pix.width, "h": pix.height})
                except Exception:
                    continue
                finally:
                    pix = None
    finally:
        doc.close()
    # Largest-area first: hero panels lead.
    out.sort(key=lambda d: d["w"] * d["h"], reverse=True)
    return out, ("fitz" if out else "fitz-no-usable-images")


def _resolve_slide_figure(ref, pool_dir, pdf_assets):
    """Map a slide's `figure` (basename or path) to a real file. Returns Path or None."""
    if not ref:
        return None
    p = Path(ref)
    if p.is_file():
        return p
    name = p.name
    if pool_dir:
        cand = pool_dir / name
        if cand.is_file():
            return cand
        # case-insensitive / suffix-tolerant match within the pool
        for f in sorted(pool_dir.iterdir()):
            if f.is_file() and f.name.lower() == name.lower():
                return f
    for a in pdf_assets:
        if a["path"].name == name:
            return a["path"]
    return None


def _gather_assets(slides, ref, quest_id, params, asset_dir):
    """Build the asset plan: per-slide resolved figure + the leftover PDF/pool 'story' figures.

    Returns (slide_figs: dict[int->Path], used: list[dict], manifest_rows: list[dict], method: str).
    """
    paper_pdf = params.get("paper_pdf") or ref.get("paper_pdf")
    figures_dir = params.get("figures_dir") or ref.get("figures_dir")
    pdf_assets, method = _extract_pdf_figures(paper_pdf, asset_dir)
    pool = _figures_root(quest_id, figures_dir)

    slide_figs, used, rows = {}, [], []
    used_paths = set()

    # 1) Honor explicit per-slide figure references first.
    for i, sl in enumerate(slides):
        fp = _resolve_slide_figure(sl.get("figure"), pool, pdf_assets)
        if fp and fp.is_file():
            slide_figs[i] = fp
            used_paths.add(str(fp.resolve()))
            rows.append({"asset": fp.name, "slide": i + 1, "title": sl.get("title", ""),
                         "src": sl.get("source", ""), "how": "explicit"})

    # 2) Auto-assign remaining "story" figures to evidence slides that lack one.
    #    Pool figures (on disk) lead over PDF-extracted; both ordered hero-first.
    story = []
    if pool:
        story += [f for f in sorted(pool.iterdir())
                  if f.is_file() and f.suffix.lower() in _IMG_EXT and str(f.resolve()) not in used_paths]
    story += [a["path"] for a in pdf_assets if str(a["path"].resolve()) not in used_paths]

    # Heuristic: skip slide 0 (cover) and conclusion-ish slides; fill the body.
    body_idx = [i for i in range(len(slides)) if i not in slide_figs and i != 0]
    si = 0
    for i in body_idx:
        if si >= len(story):
            break
        fp = story[si]; si += 1
        slide_figs[i] = fp
        used_paths.add(str(fp.resolve()))
        page = next((a["page"] for a in pdf_assets if a["path"] == fp), None)
        rows.append({"asset": fp.name, "slide": i + 1, "title": slides[i].get("title", ""),
                     "src": (f"p.{page}" if page else "figures pool"), "how": "auto"})

    for i, fp in slide_figs.items():
        used.append({"slide": i + 1, "path": str(fp)})
    return slide_figs, used, rows, method, len(pdf_assets), (str(pool) if pool else None)


def _write_asset_manifest(asset_dir, rows, method, n_pdf, pool):
    """Write asset_manifest.md (figure traceability). Returns its path (or None when no assets)."""
    if not rows and not n_pdf:
        return None
    asset_dir.mkdir(parents=True, exist_ok=True)
    # asset_manifest.md lands alongside the deck (deck dir), beside qa_report.md.
    mf = asset_dir.parent.parent / "asset_manifest.md"
    lines = ["# Asset manifest — nature-paper2ppt", "",
             f"- extraction method: `{method}`",
             f"- PDF-extracted images: {n_pdf}",
             f"- on-disk figure pool: {pool or '(none)'}",
             f"- assets placed on slides: {len(rows)}", "",
             "| asset | slide | slide title | source | how |",
             "|---|---|---|---|---|"]
    for r in rows:
        lines.append(f"| `{r['asset']}` | {r['slide']} | {_md(r['title'])} | {_md(r['src'])} | {r['how']} |")
    if not rows:
        lines.append("| _(no figures placed on slides)_ |  |  |  |  |")
    mf.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return mf


def _md(s):
    return (s or "").replace("|", "\\|").replace("\n", " ")


# ───────────────────────────── pptx authoring ─────────────────────────────
def _build_pptx(title, author, slides, slide_figs, out_path, lang):
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

    ACCENT = RGBColor(0x0A, 0x4A, 0x7A)      # restrained Nature-style blue
    INK = RGBColor(0x1A, 0x1A, 0x1A)
    MUTED = RGBColor(0x6B, 0x6B, 0x6B)
    FONT = "Microsoft YaHei" if lang == "cn" else "Calibri"

    prs = Presentation()
    prs.slide_width = Inches(13.333)         # 16:9 widescreen
    prs.slide_height = Inches(7.5)
    SW, SH = prs.slide_width, prs.slide_height
    blank = prs.slide_layouts[6]

    def _accent_bar(slide):
        bar = slide.shapes.add_shape(1, 0, 0, SW, Inches(0.12))  # MSO_SHAPE.RECTANGLE = 1
        bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()
        bar.shadow.inherit = False

    def _textbox(slide, l, t, w, h):
        tb = slide.shapes.add_textbox(l, t, w, h)
        tf = tb.text_frame
        tf.word_wrap = True
        return tf

    # ---- Cover slide ----
    cover = prs.slides.add_slide(blank)
    _accent_bar(cover)
    tf = _textbox(cover, Inches(0.9), Inches(2.2), Inches(11.5), Inches(2.6))
    p = tf.paragraphs[0]; p.text = title
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = ACCENT; p.font.name = FONT
    if author:
        p2 = tf.add_paragraph(); p2.text = author
        p2.font.size = Pt(20); p2.font.color.rgb = MUTED; p2.font.name = FONT
    cover.notes_slide.notes_text_frame.text = (slides[0].get("notes", "") if slides else "") or ""

    # ---- Body slides ----
    for i, sl in enumerate(slides):
        s = prs.slides.add_slide(blank)
        _accent_bar(s)
        fig = slide_figs.get(i)

        # Title
        ttf = _textbox(s, Inches(0.6), Inches(0.35), Inches(12.1), Inches(1.0))
        tp = ttf.paragraphs[0]; tp.text = sl.get("title", "")
        tp.font.size = Pt(28); tp.font.bold = True; tp.font.color.rgb = ACCENT; tp.font.name = FONT

        if fig:
            # Asymmetric hero layout: figure dominates (left ~62%), narrow interpretation rail (right).
            fig_l, fig_t = Inches(0.6), Inches(1.5)
            fig_box_w, fig_box_h = Inches(7.7), Inches(5.4)
            try:
                from PIL import Image
                with Image.open(str(fig)) as im:
                    iw, ih = im.size
                ar = iw / ih if ih else 1.0
                w = fig_box_w
                h = Emu(int(w / ar))
                if h > fig_box_h:
                    h = fig_box_h
                    w = Emu(int(h * ar))
                s.shapes.add_picture(str(fig), fig_l, fig_t, width=w, height=h)
            except Exception:
                s.shapes.add_picture(str(fig), fig_l, fig_t, width=fig_box_w)
            # source label under the figure
            src = sl.get("source") or ""
            if src:
                stf = _textbox(s, fig_l, Inches(6.95), fig_box_w, Inches(0.4))
                sp = stf.paragraphs[0]; sp.text = f"Source: {src}"
                sp.font.size = Pt(10); sp.font.italic = True; sp.font.color.rgb = MUTED; sp.font.name = FONT
            rail_l, rail_w = Inches(8.6), Inches(4.1)
        else:
            rail_l, rail_w = Inches(0.7), Inches(11.9)

        # Bullets + caption + takeaway in the rail / body
        btf = _textbox(s, rail_l, Inches(1.6), rail_w, Inches(5.0))
        first = True
        for b in sl.get("bullets", []):
            p = btf.paragraphs[0] if first else btf.add_paragraph()
            first = False
            p.text = "• " + b
            p.font.size = Pt(16 if fig else 20); p.font.color.rgb = INK; p.font.name = FONT
            p.space_after = Pt(8)
        cap = sl.get("caption")
        if cap:
            cp = btf.add_paragraph(); cp.text = cap
            cp.font.size = Pt(12); cp.font.italic = True; cp.font.color.rgb = MUTED; cp.font.name = FONT
            cp.space_before = Pt(8)
        take = sl.get("takeaway")
        if take:
            tp2 = btf.add_paragraph(); tp2.text = ("结论：" if lang == "cn" else "Takeaway: ") + take
            tp2.font.size = Pt(14); tp2.font.bold = True; tp2.font.color.rgb = ACCENT; tp2.font.name = FONT
            tp2.space_before = Pt(10)

        # Speaker notes
        note = sl.get("notes") or ""
        if note:
            s.notes_slide.notes_text_frame.text = note

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    return len(prs.slides)


# ───────────────────────────── QA (reopen + inspect) ─────────────────────────────
def _qa_pptx(out_path):
    """Reopen the produced .pptx, count slides / embedded media / notes. Returns a dict + writes qa_report.md."""
    report = {"slides": 0, "media": 0, "notes": 0, "ok": False, "method": "python-pptx-reopen"}
    try:
        from pptx import Presentation
        prs = Presentation(str(out_path))
        report["slides"] = len(prs.slides)
        for s in prs.slides:
            if s.has_notes_slide and (s.notes_slide.notes_text_frame.text or "").strip():
                report["notes"] += 1
        report["ok"] = report["slides"] > 0
    except Exception as e:
        report["error"] = str(e)[:200]
    # Count embedded media straight from the OOXML package (robust regardless of shape walking).
    try:
        with zipfile.ZipFile(str(out_path)) as z:
            report["media"] = sum(1 for n in z.namelist() if n.startswith("ppt/media/"))
    except Exception:
        pass
    qa = out_path.parent / "qa_report.md"
    lines = ["# QA report — nature-paper2ppt", "",
             f"- deck: `{out_path.name}`",
             f"- verification: {report['method']}",
             f"- slide count: {report['slides']}",
             f"- embedded media (ppt/media/): {report['media']}",
             f"- slides with speaker notes: {report['notes']}",
             f"- reopen OK: {report['ok']}"]
    if report.get("error"):
        lines.append(f"- reopen error: {report['error']}")
    lines += ["", "Checks performed: reopened the .pptx, counted slides, counted embedded media via the OOXML",
              "package, and verified speaker-note presence. No headless renderer was used (Lean Operating Mode)."]
    qa.write_text("\n".join(lines) + "\n", encoding="utf-8")
    return report, qa


# ───────────────────────────── HTML fallback ─────────────────────────────
def _build_html(title, author, slides, slide_figs, out_path):
    e = html.escape
    css = ("body{margin:0;font-family:Arial,Helvetica,sans-serif;background:#222}"
           "section{box-sizing:border-box;width:100%;min-height:100vh;padding:6% 8%;color:#111;"
           "background:#fff;border-top:6px solid #0a4a7a;page-break-after:always}"
           "h1{font-size:2.4em;color:#0a4a7a}h2{font-size:1.7em;color:#0a4a7a}"
           "img{max-width:60%;max-height:55vh;display:block;margin:1em 0}"
           "li{font-size:1.2em;margin:.4em 0}.cap{color:#888;font-style:italic;font-size:.9em}"
           ".take{color:#0a4a7a;font-weight:bold;margin-top:.6em}"
           ".notes{margin-top:1.2em;color:#888;font-style:italic;font-size:.9em}")
    out = [f"<!doctype html><html><head><meta charset='utf-8'><title>{e(title)}</title>"
           f"<style>{css}</style></head><body>"]
    out.append(f"<section><h1>{e(title)}</h1>" + (f"<p>{e(author)}</p>" if author else "") + "</section>")
    for i, sl in enumerate(slides):
        chunk = [f"<section><h2>{e(sl.get('title',''))}</h2>"]
        fig = slide_figs.get(i)
        if fig:
            try:
                rel = Path(fig).resolve().as_uri()
            except Exception:
                rel = str(fig)
            chunk.append(f"<img src='{e(rel)}' alt='figure'>")
        chunk.append("<ul>" + "".join(f"<li>{e(b)}</li>" for b in sl.get("bullets", [])) + "</ul>")
        if sl.get("caption"):
            chunk.append(f"<div class='cap'>{e(sl['caption'])}</div>")
        if sl.get("source"):
            chunk.append(f"<div class='cap'>Source: {e(sl['source'])}</div>")
        if sl.get("takeaway"):
            chunk.append(f"<div class='take'>结论：{e(sl['takeaway'])}</div>")
        if sl.get("notes"):
            chunk.append(f"<div class='notes'>Notes: {e(sl['notes'])}</div>")
        chunk.append("</section>")
        out.append("".join(chunk))
    out.append("</body></html>")
    out_path.parent.mkdir(parents=True, exist_ok=True)
    out_path.write_text("".join(out), encoding="utf-8")


# ───────────────────────────── entrypoint ─────────────────────────────
def render(*, command, input_path, out_path, params, quest_id):
    if not input_path:
        raise ValueError("nature-paper2ppt: --input is required (outline JSON/MD)")
    params = params or {}
    title_p, author, slides, pdf_in, figdir_in = _parse(input_path)
    title = params.get("title") or title_p
    lang = params.get("lang", "cn")
    # Merge figure-source hints (params override file).
    src_hints = {"paper_pdf": params.get("paper_pdf") or pdf_in,
                 "figures_dir": params.get("figures_dir") or figdir_in}

    out_path = Path(out_path)
    out_path.parent.mkdir(parents=True, exist_ok=True)
    asset_dir = out_path.parent / "assets" / "figures"

    # Gather figure assets (PDF extraction degrades gracefully; pool used regardless).
    slide_figs, used, rows, method, n_pdf, pool = _gather_assets(
        slides, src_hints, quest_id, {**params, **src_hints}, asset_dir)
    manifest = _write_asset_manifest(asset_dir, rows, method, n_pdf, pool)

    # ---- PRIMARY: real .pptx via python-pptx ----
    try:
        import pptx  # noqa: F401
        have_pptx = True
    except Exception:
        have_pptx = False

    if have_pptx:
        pptx_path = out_path if out_path.suffix.lower() == ".pptx" else out_path.with_suffix(".pptx")
        n_slides = _build_pptx(title, author, slides, slide_figs, pptx_path, lang)
        qa, qa_path = _qa_pptx(pptx_path)
        return {"ok": True, "out_path": str(pptx_path), "format": "pptx",
                "summary": (f"real .pptx: {n_slides} slides, {qa['media']} media, "
                            f"{qa['notes']} notes (figs:{method})"),
                "meta": {"slides": n_slides, "media": qa["media"], "notes": qa["notes"],
                         "title": title, "lang": lang, "fig_method": method,
                         "pdf_images": n_pdf, "figures_used": used,
                         "qa_report": str(qa_path),
                         "asset_manifest": str(manifest) if manifest else None,
                         "renderer": "python-pptx"}}

    # ---- FALLBACK: self-contained HTML deck (python-pptx unavailable) ----
    html_path = out_path if out_path.suffix.lower() == ".html" else out_path.with_suffix(".html")
    _build_html(title, author, slides, slide_figs, html_path)
    return {"ok": True, "out_path": str(html_path), "format": "html",
            "summary": (f"python-pptx unavailable; HTML deck fallback: {len(slides) + 1} slides "
                        f"(figs:{method})"),
            "meta": {"slides": len(slides) + 1, "title": title, "lang": lang, "fig_method": method,
                     "pdf_images": n_pdf, "figures_used": used,
                     "asset_manifest": str(manifest) if manifest else None,
                     "renderer": "html-fallback",
                     "note": "install python-pptx to emit a real .pptx at this ref (no contract change)"}}
